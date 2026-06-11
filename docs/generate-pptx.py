#!/usr/bin/env python3
"""Generate Operiq AI Presentation (PPTX) using python-pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

ACCENT = RGBColor(0x10, 0xA3, 0x7F)  # #10a37f
DARK = RGBColor(0x18, 0x18, 0x1B)     # zinc-900
BODY = RGBColor(0x3F, 0x3F, 0x46)     # zinc-600
LIGHT_BG = RGBColor(0xF4, 0xF4, 0xF5) # zinc-100

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide(bg_color=RGBColor(0xFF, 0xFF, 0xFF)):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    return slide

def add_title(slide, text, left=0.8, top=0.5, width=11.7, size=32, color=DARK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    return tf

def add_subtitle(slide, text, left=0.8, top=1.5, width=11.7, size=18, color=BODY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    return tf

def add_body(slide, text, left=0.8, top=2.3, width=11.7, size=16, color=BODY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.space_after = Pt(8)
    return tf

def add_bullets(slide, items, left=1.2, top=2.3, width=11, size=16, color=BODY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.level = 0
    return tf

def add_accent_bar(slide, left=0, top=0, height=0.06):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(left), Inches(top), Inches(13.333), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

# ====== Slide 1: Title ======
slide1 = add_slide()
add_accent_bar(slide1, top=0)
add_title(slide1, "Operiq AI", left=0.8, top=2.2, size=48, color=ACCENT)
add_subtitle(slide1, "Workplace Productivity Assistant", left=0.8, top=3.2, size=28)
add_body(slide1, "AI-powered platform for professionals: email drafting, meeting summaries, task planning, research, and intelligent chat.", left=0.8, top=4.2, size=16, color=RGBColor(0x71, 0x71, 0x7A))
# Watermark attempt - add small logo text
txBox = slide1.shapes.add_textbox(Inches(10), Inches(6.5), Inches(3), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "v2.0.0"
p.font.size = Pt(10)
p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
p.alignment = PP_ALIGN.RIGHT

# ====== Slide 2: Problem Statement ======
slide2 = add_slide()
add_accent_bar(slide2, top=0)
add_title(slide2, "Problem Statement", size=30)
add_bullets(slide2, [
    "Professionals spend up to 60% of their time on administrative and communication tasks",
    "Repetitive activities: drafting emails, summarizing meetings, planning schedules, researching",
    "These tasks divert attention from strategic work and creative problem-solving",
    "Result: reduced productivity, delayed decision-making, professional burnout",
    "",
    "Operiq AI solves this by automating and accelerating repetitive workplace tasks",
    "Leveraging large language models and intelligent automation",
], top=1.8)

# ====== Slide 3: Solution Overview ======
slide3 = add_slide()
add_accent_bar(slide3, top=0)
add_title(slide3, "Solution Overview", size=30)
add_bullets(slide3, [
    "Smart Email Generator - Professional emails with tone/audience adaptation",
    "Meeting Notes Summarizer - Key points, decisions, action items extraction",
    "AI Task Planner - Structured daily/weekly plans with prioritization",
    "AI Research Assistant - Article summarization with key insights",
    "AI Chatbot Interface - Interactive multi-turn workplace assistant",
    "",
    "All powered by advanced AI models with prompt engineering and responsible AI practices",
], top=1.8)

# ====== Slide 4: System Architecture ======
slide4 = add_slide()
add_accent_bar(slide4, top=0)
add_title(slide4, "System Architecture", size=30)
add_bullets(slide4, [
    "Netlify CDN + Edge Functions - Global deployment, request routing, security headers",
    "Convex Serverless Backend - Real-time database, ACID transactions, built-in auth",
    "AI Providers - OpenAI-compatible gateway, Hugging Face (images/video), ElevenLabs (voice)",
    "External Services - Resend (email), DuckDuckGo (web search)",
    "",
    "Architecture diagram available in documentation",
], top=1.8)
# Try to embed the architecture diagram
try:
    img_path = "docs/operiq-system-architecture.png"
    if os.path.exists(img_path):
        slide4.shapes.add_picture(img_path, Inches(0.5), Inches(4.2), Inches(12.3), height=Inches(3))
except:
    pass

# ====== Slide 5: Software Architecture ======
slide5 = add_slide()
add_accent_bar(slide5, top=0)
add_title(slide5, "Software Architecture", size=30)
add_bullets(slide5, [
    "Frontend: React 19 + TanStack Start + Tailwind CSS v4 + shadcn/ui + lucide-react",
    "State: Convex useQuery (real-time reactive) + React useState (UI) + @convex-dev/auth",
    "Backend: Convex modules - schema, users, emailDrafts, summaries, plans, analyses, billing",
    "Tool System: Registry-based with web_search, fetch_url, generate_image, generate_video",
    "Agent Mode: ReAct loop with multi-step reasoning and tool execution",
], top=1.8)

# ====== Slide 6: AI Pipeline ======
slide6 = add_slide()
add_accent_bar(slide6, top=0)
add_title(slide6, "AI Pipeline", size=30)
add_bullets(slide6, [
    "Model Routing: Task type determines model and provider selection",
    "Prompt Assembly: System instructions + user context + output format + guardrails",
    "Generation: Streaming responses with real-time display",
    "Post-processing: Validation, markdown rendering, disclaimer injection",
    "",
    "Provider Mix: OpenAI-compatible gateway, Hugging Face, ElevenLabs",
    "Tool Integration: Agents can search web, fetch URLs, generate images/video",
], top=1.8)

# ====== Slide 7: Prompt Engineering ======
slide7 = add_slide()
add_accent_bar(slide7, top=0)
add_title(slide7, "Prompt Engineering", size=30)
add_bullets(slide7, [
    "System Instructions: Define persona, output format, boundaries, and tone",
    "Context Injection: User-provided data clearly delineated from system instructions",
    "Output Guardrails: Structured formats, validation checks, anti-hallucination rules",
    "",
    "Email Prompt: Professional tone, audience-aware, no fabricated information",
    "Meeting Prompt: Extract decisions, action items with assignees, deadlines",
    "Planner Prompt: Time-optimized schedules with Eisenhower Matrix prioritization",
    "Research Prompt: Key insights extraction, complexity simplification",
], top=1.8)

# ====== Slide 8: Responsible AI ======
slide8 = add_slide()
add_accent_bar(slide8, top=0)
add_title(slide8, "Responsible AI", size=30)
add_bullets(slide8, [
    "Bias Mitigation - Neutrality instructions in all prompts, output monitoring",
    "Data Privacy - AI providers configured to never train on user data",
    "Disclaimers - Every AI output includes review reminder and AI-generated label",
    "Rate Limiting - Free tier: 50 AI requests/day; Pro/Enterprise: higher limits",
    "Transparency - Clear labeling of AI-generated content throughout the platform",
    "DevTools Protection - Production-only guard against casual code inspection",
], top=1.8)

# ====== Slide 9: Evaluation Criteria ======
slide9 = add_slide()
add_accent_bar(slide9, top=0)
add_title(slide9, "Evaluation Criteria", size=30)
add_bullets(slide9, [
    "Problem Relevance (20%) - Addresses real workplace productivity challenges",
    "Prompt Engineering (25%) - Structured, tested, refined prompts with guardrails",
    "Functionality Accuracy (25%) - AI-powered features with validation checks",
    "Innovation (15%) - Multi-agent loop, voice TTS/STT, image/video generation",
    "Responsible AI (10%) - Bias detection, privacy safeguards, disclaimers",
    "Presentation Clarity (5%) - Professional documentation with architecture diagrams",
], top=1.8)

# ====== Slide 10: Tools & Technologies ======
slide10 = add_slide()
add_accent_bar(slide10, top=0)
add_title(slide10, "Tools & Technologies", size=30)
add_bullets(slide10, [
    "AI/ML: OpenAI API, Hugging Face (FLUX, HunyuanVideo), ElevenLabs (TTS/STT)",
    "Frontend: React 19, TanStack Start/Router, Tailwind CSS v4, shadcn/ui, lucide-react",
    "Backend: Convex (real-time DB, auth, file storage), TypeScript",
    "Infrastructure: Netlify (CDN, Edge Functions, serverless), Vite, Bun",
    "External: Resend (email delivery), DuckDuckGo (web search), Stripe (demo billing)",
    "Tool System: Registry-based agent tools (web_search, fetch_url, generate_image)",
], top=1.8)

# ====== Slide 11: Sample Prompts ======
slide11 = add_slide()
add_accent_bar(slide11, top=0)
add_title(slide11, "Sample Prompts & Outputs", size=30)
add_bullets(slide11, [
    "Email: Input 'Update to John (Manager), Professional tone'",
    "  Output: Subject line, greeting, status update, next steps, closing",
    "",
    "Meeting Summary: Input '30-min team sync: Q2 roadmap, hiring, sprint retro'",
    "  Output: Key Decisions, Action Items (with assignees), Discussion Highlights",
    "",
    "Research: Input 'Summarize market trends in AI productivity tools 2025'",
    "  Output: Executive summary, key trends, competitive landscape, recommendations",
], top=1.8)

# ====== Slide 12: Challenges & Solutions ======
slide12 = add_slide()
add_accent_bar(slide12, top=0)
add_title(slide12, "Challenges & Solutions", size=30)
add_bullets(slide12, [
    "SSR Auth: TanStack Start SSR + Convex auth. Solved with ConvexAuthProvider + AuthGate",
    "AI Streaming: AI SDK streaming with custom tools. Solved with streamText + maxSteps",
    "Netlify Deploy: Env vars in Workers. Solved with .server.ts request-time pattern",
    "Multi-Provider: Multiple AI providers. Solved with unified interface + model mapping",
    "Usage Tracking: Real-time without blocking. Solved with fire-and-forget mutations",
    "Billing Reset: Period-based limits. Solved with auto-reset on period expiry check",
], top=1.8)

# ====== Slide 13: Thank You ======
slide13 = add_slide()
add_accent_bar(slide13, top=0)
add_title(slide13, "Thank You", left=0.8, top=2.5, size=48, color=ACCENT)
add_subtitle(slide13, "Operiq AI - Workplace Productivity, Reimagined", left=0.8, top=3.5, size=24)
add_body(slide13, "operiq.ai", left=0.8, top=4.5, size=18, color=RGBColor(0x71, 0x71, 0x7A))

# Save
output_path = "docs/operiq-presentation.pptx"
prs.save(output_path)
print(f"PPTX generated: {output_path} ({len(prs.slides)} slides)")
